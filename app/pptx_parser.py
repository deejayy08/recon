# app/pptx_parser.py
from pptx import Presentation
import uuid

def extract_chunks_from_pptx(path: str, s3_uri: str):
    prs = Presentation(path)
    chunks = []
    slide_idx = 0
    for slide in prs.slides:
        slide_idx += 1
        # text from shapes
        text_blocks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_blocks.append(shape.text.strip())
            # tables
            if shape.has_table:
                table = shape.table
                # table rows
                for r, row in enumerate(table.rows, start=1):
                    cells = []
                    for c, cell in enumerate(row.cells, start=1):
                        txt = cell.text.strip()
                        cells.append(txt)
                        chunk_cell = {
                            "chunk_id": uuid.uuid4().hex,
                            "text": txt,
                            "metadata": {"doc_uri": s3_uri, "slide": slide_idx, "table": True, "row": r, "col": c}
                        }
                        chunks.append(chunk_cell)
                    # add row-level chunk
                    chunk_row = {
                        "chunk_id": uuid.uuid4().hex,
                        "text": " | ".join(cells),
                        "metadata": {"doc_uri": s3_uri, "slide": slide_idx, "table": True, "row": r}
                    }
                    chunks.append(chunk_row)
        # add slide-level text chunk
        if text_blocks:
            chunk = {"chunk_id": uuid.uuid4().hex, "text": "\n".join(text_blocks), "metadata": {"doc_uri": s3_uri, "slide": slide_idx}}
            chunks.append(chunk)
    return chunks
